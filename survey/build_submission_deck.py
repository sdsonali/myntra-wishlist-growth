"""
10-slide NextLeap submission deck (16:9, 14pt everywhere, no personal name).

  python survey/build_submission_deck.py
"""

from __future__ import annotations

from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "Myntra_Wishlist_Confidence_Gap_SUBMISSION.pptx"
ASSETS = ROOT / "deck" / "assets"
CHART = ASSETS / "blockers.png"

APP = "https://myntra-wishlist-growth.streamlit.app/"
FORM = (
    "https://docs.google.com/forms/d/e/"
    "1FAIpQLSdPOnDtw1lQGARNYOo2zOoUMpwlOAcLxIzmPaLz27XYN0TjNQ/viewform"
)
AMAZON = "https://www.amazon.in/fashion"
FLIPKART = "https://www.flipkart.com/"
AJIO = "https://www.ajio.com/"
NYKAA = "https://www.nykaafashion.com/"

NAVY = RGBColor(0x1E, 0x3A, 0x5F)
NAVY2 = RGBColor(0x0D, 0x24, 0x40)
AMBER = RGBColor(0xC4, 0x6B, 0x2F)
AMBER_LT = RGBColor(0xF6, 0xEB, 0xDC)
COOL = RGBColor(0xE6, 0xEE, 0xF4)
WARM = RGBColor(0xF8, 0xEE, 0xE4)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x4A, 0x55, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC8, 0xC4, 0xBC)
BG = RGBColor(0xF7, 0xF5, 0xF1)
OUT_GREY = RGBColor(0x6B, 0x72, 0x78)
LINK = RGBColor(0x1F, 0x4E, 0x79)
PINK = RGBColor(0xFF, 0x3F, 0x6C)
OK = RGBColor(0x2C, 0x3E, 0x50)
SHIP = RGBColor(0x1E, 0x3A, 0x5F)
PLAN = RGBColor(0xC4, 0x6B, 0x2F)

W = Inches(13.333)
H = Inches(7.5)
FS = Pt(14)
FONT = "Calibri"


def I(n: float):
    return Inches(n)


def _run(run, text, *, bold=False, color=INK, italic=False, underline=False, url=None):
    run.text = text
    run.font.name = FONT
    run.font.size = FS
    run.font.bold = bold
    run.font.italic = italic
    run.font.underline = underline or bool(url)
    run.font.color.rgb = LINK if url else color
    if url:
        run.hyperlink.address = url


def _fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()
    shape.shadow.inherit = False


def rect(slide, l, t, w, h, fill, *, line=None, rounded=False, adj=0.08):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        I(l),
        I(t),
        I(w),
        I(h),
    )
    if rounded:
        sh.adjustments[0] = adj
    _fill(sh, fill)
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    return sh


def tf_shape(shape, *, anchor=MSO_ANCHOR.TOP, ml=0.08, mr=0.08, mt=0.04, mb=0.04):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = I(ml)
    tf.margin_right = I(mr)
    tf.margin_top = I(mt)
    tf.margin_bottom = I(mb)
    anchor_map = {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}
    try:
        tf._txBody.bodyPr.set("anchor", anchor_map[anchor])
    except Exception:
        pass
    return tf


def _write_paras(tf, paras, default_align=PP_ALIGN.LEFT):
    first = True
    for item in paras:
        if isinstance(item, str):
            runs, opts = [(item, {})], {}
        elif isinstance(item, tuple) and item and isinstance(item[0], str):
            runs, opts = [(item[0], item[1] if len(item) > 1 else {})], item[1] if len(item) > 1 else {}
        else:
            runs, opts = item, {}
            if runs and isinstance(runs, tuple):
                runs, opts = runs
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = opts.get("align", default_align)
        p.space_after = Pt(opts.get("space_after", 2))
        p.space_before = Pt(opts.get("space_before", 0))
        p.level = 0
        for text, ro in runs if isinstance(runs, list) else [(runs, opts)]:
            if isinstance(text, tuple):
                text, ro = text
            run = p.add_run()
            _run(
                run,
                text,
                bold=ro.get("bold", False),
                color=ro.get("color", INK),
                italic=ro.get("italic", False),
                underline=ro.get("underline", False),
                url=ro.get("url"),
            )


def tb(slide, l, t, w, h, paras, *, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT, ml=0.06, mr=0.06, mt=0.02, mb=0.02):
    sh = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
    tf = tf_shape(sh, anchor=anchor, ml=ml, mr=mr, mt=mt, mb=mb)
    if isinstance(paras, (str, tuple)):
        paras = [paras]
    _write_paras(tf, paras, default_align=align)
    return sh


def card(slide, l, t, w, h, fill=WHITE, line=LINE, rounded=True):
    return rect(slide, l, t, w, h, fill, line=line, rounded=rounded)


def chip(slide, l, t, w, h, label, fill, color=WHITE):
    rect(slide, l, t, w, h, fill, rounded=True, adj=0.2)
    tb(
        slide,
        l,
        t,
        w,
        h,
        [(label, {"bold": True, "color": color, "align": PP_ALIGN.CENTER})],
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.CENTER,
    )


def footer(slide, n):
    rect(slide, 0, 7.22, 13.333, 0.28, NAVY)
    tb(
        slide,
        0.28,
        7.22,
        10.2,
        0.28,
        [("Product Manager, Growth Team", {"bold": True, "color": WHITE})],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    tb(
        slide,
        11.4,
        7.22,
        1.65,
        0.28,
        [(str(n), {"bold": True, "color": WHITE, "align": PP_ALIGN.RIGHT})],
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.RIGHT,
    )


def header(slide, eyebrow, title):
    rect(slide, 0, 0, 13.333, 7.5, BG)
    rect(slide, 0, 0, 0.12, 7.5, NAVY)
    tb(slide, 0.32, 0.08, 12.7, 0.28, [(eyebrow, {"bold": True, "color": AMBER})])
    tb(slide, 0.32, 0.32, 12.7, 0.55, [(title, {"bold": True, "color": NAVY})])
    rect(slide, 0.32, 0.90, 12.7, 0.015, AMBER)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_cell(cell, text, *, bold=False, color=INK, fill=None, align=PP_ALIGN.LEFT, url=None):
    cell.text = ""
    cell.margin_left = I(0.05)
    cell.margin_right = I(0.05)
    cell.margin_top = I(0.03)
    cell.margin_bottom = I(0.03)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    _run(run, text, bold=bold, color=color, url=url)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.background()
    for border in ("lnL", "lnR", "lnT", "lnB"):
        ln = cell._tc.get_or_add_tcPr().find(qn(f"a:{border}"))
        if ln is None:
            continue


def table(slide, l, t, w, h, rows, col_w, header=True):
    """rows: list of list of str | (str, dict). col_w: relative weights."""
    n_rows, n_cols = len(rows), len(rows[0])
    shp = slide.shapes.add_table(n_rows, n_cols, I(l), I(t), I(w), I(h))
    tbl = shp.table
    total = sum(col_w)
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = I(w * cw / total)
    for r, row in enumerate(rows):
        for c, item in enumerate(row):
            if isinstance(item, str):
                text, opts = item, {}
            else:
                text, opts = item[0], item[1]
            is_head = header and r == 0
            set_cell(
                tbl.cell(r, c),
                text,
                bold=opts.get("bold", is_head),
                color=opts.get("color", WHITE if is_head else INK),
                fill=opts.get("fill", NAVY if is_head else (COOL if r % 2 == 0 else WHITE)),
                align=opts.get("align", PP_ALIGN.LEFT),
                url=opts.get("url"),
            )
    return shp


def make_blocker_chart(path: Path):
    path.parent.mkdir(parents=True, exist_ok=True)
    labels = [
        "Quality doubt",
        "Price change",
        "Fit uncertainty",
        "Found alternative",
        "No urgency",
    ]
    pcts = [44.7, 23.8, 20.3, 7.9, 2.1]
    scopes = ["IN", "OUT", "IN", "IN", "IN"]
    colors = ["#1E3A5F" if s == "IN" else "#6B7278" for s in scopes]
    fig, ax = plt.subplots(figsize=(6.15, 2.85), dpi=160)
    fig.patch.set_facecolor("#F7F5F1")
    ax.set_facecolor("#F7F5F1")
    y = list(range(len(labels)))[::-1]
    ax.barh(y, pcts[::-1], color=colors[::-1], height=0.62, zorder=2)
    for yi, pct, scope, lab in zip(y, pcts[::-1], scopes[::-1], labels[::-1]):
        ax.text(
            pct + 0.8,
            yi,
            f"{pct:.1f}%   {scope}",
            va="center",
            ha="left",
            fontsize=14,
            fontname="Calibri",
            color="#1A1A1A",
            fontweight="bold",
        )
        ax.text(-0.6, yi, lab, va="center", ha="right", fontsize=14, fontname="Calibri", color="#1E3A5F")
    ax.set_xlim(0, 62)
    ax.set_yticks([])
    ax.set_xlabel("% of 340 purchase-blocker mentions", fontsize=14, fontname="Calibri", color="#4A555C")
    ax.tick_params(axis="x", labelsize=14)
    for lbl in ax.get_xticklabels():
        lbl.set_fontname("Calibri")
        lbl.set_fontsize(14)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_visible(False)
    ax.grid(axis="x", color="#C8C4BC", linewidth=0.6, zorder=0)
    fig.tight_layout(pad=0.3)
    fig.savefig(path, dpi=160, bbox_inches="tight", facecolor=fig.get_facecolor())
    plt.close(fig)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------


def slide_1(prs):
    s = blank(prs)
    rect(s, 0, 0, 13.333, 7.5, NAVY)
    rect(s, 0, 6.55, 13.333, 0.95, AMBER)
    tb(s, 0.55, 0.55, 12.2, 0.35, [("PRODUCT STRATEGY  —  MYNTRA GROWTH", {"bold": True, "color": AMBER})])
    tb(
        s,
        0.55,
        1.05,
        12.2,
        1.15,
        [("CLOSING THE WISHLIST CONFIDENCE GAP", {"bold": True, "color": WHITE})],
    )
    tb(
        s,
        0.55,
        2.25,
        12.0,
        0.85,
        [
            (
                "Improving Myntra's wishlist-to-purchase conversion by resolving in-app decision doubt — without discounts, coupons, or cashback.",
                {"color": WHITE},
            )
        ],
    )
    card(s, 0.55, 3.25, 12.2, 0.85, fill=RGBColor(0x16, 0x2E, 0x4C), line=AMBER)
    tb(
        s,
        0.7,
        3.35,
        11.9,
        0.7,
        [
            (
                "For 24–30 shoppers with at least one unsold saved item, resolve fit / quality / compare doubt on the wishlist — so WPCR moves without discounts.",
                {"color": WHITE},
            )
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    chips = [
        ("FOCUS", "Ethnic / occasion / fashion-forward saved items"),
        ("LEVER", "In-app decision confidence — not price"),
        ("CONSTRAINT", "No monetary incentives of any kind"),
    ]
    for i, (k, v) in enumerate(chips):
        x = 0.55 + i * 4.1
        card(s, x, 4.3, 3.9, 1.85, fill=WHITE, line=None)
        tb(s, x + 0.15, 4.42, 3.6, 0.35, [(k, {"bold": True, "color": AMBER})])
        tb(s, x + 0.15, 4.78, 3.6, 1.15, [(v, {"color": NAVY, "bold": True})])
    tb(
        s,
        0.55,
        6.62,
        10.4,
        0.75,
        [
            [
                ("Resources  ·  ", {"bold": True, "color": WHITE}),
                ("Live app (Discovery + MVP)", {"url": APP, "bold": True}),
                ("   ·   ", {"color": WHITE}),
                ("Interview form", {"url": FORM, "bold": True}),
                ("   ·   Product Manager, Growth Team", {"bold": True, "color": WHITE}),
            ]
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    tb(
        s,
        11.6,
        6.62,
        1.4,
        0.75,
        [("1", {"bold": True, "color": WHITE, "align": PP_ALIGN.RIGHT})],
        anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.RIGHT,
    )


def slide_2(prs):
    s = blank(prs)
    header(s, "THE BUSINESS METRIC, BROKEN DOWN", "Saved demand dies at decision confidence, not at re-engagement")
    tb(
        s,
        0.32,
        0.96,
        12.7,
        0.32,
        [
            (
                "NORTH STAR — WPCR: share of users who purchase ≥1 wishlisted item within 30 days of adding it.",
                {"color": MUTED},
            )
        ],
    )
    stages = [
        ("Add to wishlist", False),
        ("Reconsideration trigger", False),
        ("Confidence resolution", True),
        ("Decision", False),
        ("Checkout", False),
    ]
    for i, (name, hot) in enumerate(stages):
        x = 0.32 + i * 2.58
        card(s, x, 1.32, 2.42, 0.78, fill=AMBER if hot else WHITE, line=NAVY if hot else LINE)
        tb(
            s,
            x,
            1.32,
            2.42,
            0.50,
            [(name, {"bold": True, "color": WHITE if hot else NAVY, "align": PP_ALIGN.CENTER})],
            anchor=MSO_ANCHOR.MIDDLE,
            align=PP_ALIGN.CENTER,
        )
        if hot:
            tb(
                s,
                x,
                1.72,
                2.42,
                0.32,
                [("TARGET STAGE", {"bold": True, "color": WHITE, "align": PP_ALIGN.CENTER})],
                align=PP_ALIGN.CENTER,
            )
        if i < 4:
            tb(s, x + 2.28, 1.48, 0.28, 0.32, [("→", {"bold": True, "color": AMBER})])
    tb(
        s,
        0.32,
        2.18,
        12.7,
        0.32,
        [
            (
                "Highest-potential areas from the discovery corpus, and how the MVP can move WPCR (mechanism only — no invented lift %).",
                {"bold": True, "color": NAVY},
            )
        ],
    )
    rows = [
        ["Opportunity (corpus)", "MVP surface", "How it can move WPCR"],
        ["Quality doubt 44.7%  IN", "Still deciding — this SKU's reviews", "Turns photo-vs-real freeze into a yes/no without leaving the app"],
        ["Fit 20.3%  IN", "Size + Size up, then buy / Need more info", "Cuts the “order two sizes and wait” delay"],
        ["Found alt 7.9% + leave-app 10.8%  IN", "Compare 2–3 + in-app answer", "Closes the loop on Myntra instead of another app"],
        [
            ("Price 23.8%  OUT", {"bold": True, "color": MUTED, "fill": AMBER_LT}),
            ("Out of scope", {"bold": True, "color": MUTED, "fill": AMBER_LT}),
            ("No monetary lever in this brief", {"color": MUTED, "fill": AMBER_LT}),
        ],
        ["No urgency 2.1%  IN (secondary)", "Express 30-min demo", "Optional urgency nudge; not the hero; not live logistics"],
    ]
    table(s, 0.32, 2.50, 12.7, 3.55, rows, [3.4, 3.5, 5.8])
    card(s, 0.32, 6.15, 6.2, 0.95, fill=COOL, line=NAVY)
    card(s, 6.72, 6.15, 6.3, 0.95, fill=AMBER_LT, line=OUT_GREY)
    tb(s, 0.42, 6.18, 6.0, 0.28, [("IN SCOPE", {"bold": True, "color": NAVY})])
    tb(
        s,
        0.42,
        6.44,
        6.0,
        0.60,
        [
            (
                "Confidence-blocked intenders — unsold saved items, “Probably” intent. Ethnic / occasion / fashion-forward over basics.",
                {"color": INK},
            )
        ],
    )
    tb(s, 6.82, 6.18, 6.1, 0.28, [("OUT OF SCOPE", {"bold": True, "color": MUTED})])
    tb(
        s,
        6.82,
        6.44,
        6.1,
        0.60,
        [
            (
                "Price-only waiters — no monetary lever exists. Rare / infrequent shoppers outside the monthly+ cadence.",
                {"color": INK},
            )
        ],
    )
    footer(s, 2)


def slide_3(prs):
    s = blank(prs)
    header(s, "HOW THE AI DISCOVERY ENGINE WORKS", "The engine ranks opportunities for WPCR — quality leads, price is out")
    tb(
        s,
        0.32,
        0.94,
        12.7,
        0.30,
        [
            [
                ("Job: identify, quantify, and compare opportunity areas that could move WPCR at stage 3.  ", {"color": MUTED}),
                ("Open the live engine", {"url": APP, "bold": True}),
            ]
        ],
    )
    # Left column
    card(s, 0.32, 1.28, 6.35, 5.78, fill=WHITE, line=LINE)
    tb(s, 0.44, 1.34, 6.1, 0.30, [("CLICK PATH — ASK A GROUNDED QUESTION", {"bold": True, "color": AMBER})])
    clicks = [
        "1. Open the live app → tab Discovery Engine.",
        "2. Click a theme shortcut (Fit, Quality, Price, Occasion, Comparison, Left the app) or type a PM question.",
        "3. Click Get grounded answer.",
        "4. Read numbered findings + “Validate in interviews”. Unsupported writer figures are discarded.",
        "5. Right rail Opportunity comparison: headline, 424 / 7.5% / 10.8%, source-mix, blocker and reason bars.",
    ]
    tb(s, 0.44, 1.64, 6.1, 2.15, [(c, {}) for c in clicks])
    tb(s, 0.44, 3.80, 6.1, 0.28, [("OFFLINE PIPELINE  ·  LIVE AGENT", {"bold": True, "color": NAVY})])
    tb(
        s,
        0.44,
        4.08,
        6.1,
        1.15,
        [
            (
                "Offline: scrape Play Store + App Store + YouTube → LLM curates gold → tag 4 dimensions (reason, blocker, comparison, outside research) → aggregate the opportunity table.",
                {},
            ),
            (
                "Live: Router maps the question onto the real tag vocabulary (no keyword lists) → Retrieval ranks matched signals, IDF-weighted, spread across sources → Writer drafts findings under an anti-hallucination check.",
                {},
            ),
        ],
    )
    tb(s, 0.44, 5.28, 6.1, 0.28, [("HOW THE NUMBERS ARE MADE", {"bold": True, "color": NAVY})])
    tb(
        s,
        0.44,
        5.56,
        6.1,
        1.40,
        [
            ("424 = rows in the gold tagged JSON after scrape → curate → tag. Play 154 · YouTube 212 · App Store 58. No Reddit.", {}),
            ("Blocker % = label count ÷ 340 blocker mentions (not ÷ 424 reviews).", {}),
            ("7.5% comparison = 32 / 424 rows tagged comparison-true. 10.8% external = 46 / 424 rows that mention leaving the app.", {}),
        ],
    )
    # Right column
    card(s, 6.80, 1.28, 6.22, 5.78, fill=WHITE, line=LINE)
    tb(s, 6.92, 1.34, 6.0, 0.30, [("WHAT IT FOUND — RANKED FOR WPCR", {"bold": True, "color": AMBER})])
    tb(
        s,
        6.92,
        1.64,
        6.0,
        0.55,
        [
            (
                "Quality doubt 44.7%; fit 20.3%. Price is larger at 23.8% but OUT — no monetary lever. Comparison 7.5%; leave-app 10.8%.",
                {"bold": True, "color": NAVY},
            )
        ],
    )
    s.shapes.add_picture(str(CHART), I(6.95), I(2.22), I(5.95), I(2.55))
    tb(s, 6.92, 4.80, 6.0, 0.28, [("PATTERNS (PARAPHRASED — NEVER QUOTED)", {"bold": True, "color": NAVY})])
    tb(
        s,
        6.92,
        5.08,
        6.0,
        1.15,
        [
            ("Quality: fabric/finish cheaper than the listing, or authenticity doubt.", {}),
            ("Fit: size chart is not trusted without a planned exchange.", {}),
            ("Found alt: same item checked on another app — other apps are 60.9% of outside-app mentions (YouTube 17.4% · friends 10.9% · Google 8.7%).", {}),
        ],
    )
    tb(
        s,
        6.92,
        6.22,
        6.0,
        0.72,
        [
            (
                "Caveats: public reviews under-index the silent freeze (nobody reviews an item they never bought). Quality bucket is broad — keyword fallback on null tags can absorb service/delivery noise.",
                {"italic": True, "color": MUTED},
            )
        ],
    )
    footer(s, 3)


def slide_4(prs):
    s = blank(prs)
    header(s, "THE MVP — SHOPPER UI", "On the wishlist: ask, compare, or (demo) get it in 30 minutes")
    tb(
        s,
        0.32,
        0.94,
        12.7,
        0.30,
        [
            [
                ("Same URL, tab Fit & Confidence Assistant. Grounded in this product's reviews — not the 424 discovery corpus.  ", {"color": MUTED}),
                ("Open deployed MVP", {"url": APP, "bold": True}),
            ]
        ],
    )
    # Wishlist mock
    card(s, 0.32, 1.28, 6.35, 3.05, fill=WHITE, line=LINE)
    tb(s, 0.44, 1.32, 6.1, 0.28, [("1. WISHLIST + FILTERS", {"bold": True, "color": AMBER})])
    tb(
        s,
        0.44,
        1.58,
        6.1,
        0.40,
        [
            (
                "Cards: image, brand, price, review badge. Pills All / Ethnic Wear / In stock narrow the frozen SKU. Select a card to ask about it.",
                {},
            )
        ],
    )
    for i, lab in enumerate(["All", "Ethnic Wear", "In stock"]):
        chip(s, 0.50 + i * 1.55, 2.02, 1.42, 0.32, lab, NAVY if i == 1 else COOL, WHITE if i == 1 else NAVY)
    products = [
        ("Libas", "Embroidered Anarkali", "Rs. 2499", "Reviews: enough"),
        ("Sassafras", "Sequinned Gown", "Rs. 3299", "Reviews: enough"),
        ("Ahalyaa", "Kurta + Palazzo", "Rs. 1899", "Reviews: enough"),
    ]
    for i, (brand, name, price, badge) in enumerate(products):
        x = 0.48 + i * 2.02
        card(s, x, 2.42, 1.92, 1.78, fill=COOL, line=LINE)
        rect(s, x, 2.42, 1.92, 0.08, PINK)
        tb(s, x + 0.06, 2.52, 1.80, 0.28, [(brand, {"bold": True, "color": NAVY})])
        tb(s, x + 0.06, 2.78, 1.80, 0.45, [(name, {"color": MUTED})])
        tb(s, x + 0.06, 3.20, 1.80, 0.28, [(price, {"bold": True, "color": INK})])
        chip(s, x + 0.08, 3.52, 1.72, 0.28, badge, OK)
    # Still deciding
    card(s, 6.80, 1.28, 6.22, 3.05, fill=WHITE, line=LINE)
    tb(s, 6.92, 1.32, 6.0, 0.28, [("2. STILL DECIDING", {"bold": True, "color": AMBER})])
    tb(
        s,
        6.92,
        1.60,
        6.0,
        0.70,
        [
            (
                "Usual size + occasion + free-text question → Get answer → exactly one verdict + proof bullets from this SKU's reviews → Add to bag or Still not sure.",
                {},
            )
        ],
    )
    card(s, 7.02, 2.35, 5.88, 1.82, fill=WARM, line=AMBER)
    tb(s, 7.12, 2.40, 5.7, 0.28, [("WORKED EXAMPLE  ·  RUNS SMALL", {"bold": True, "color": AMBER})])
    tb(
        s,
        7.12,
        2.68,
        5.7,
        1.40,
        [
            [("Verdict: ", {"bold": True}), ("Size up, then buy", {"bold": True, "color": NAVY})],
            ("“Go one size up from your usual — this runs a little small.”", {"italic": True}),
            ("Proof: several reviewers say it runs small and recommend ordering one size up.", {}),
        ],
    )
    # Four bottom cards
    feats = [
        (
            "3. COMPARE 2–3",
            "Tick 2–3 saved items → recommendation + table. Hard cap at 3. Breaks gown-vs-kurta freeze without opening another app.",
            COOL,
        ),
        (
            "4. 30-MIN EXPRESS (DEMO)",
            "Pincode → Check. Hub + size in stock → Get it in 30 mins. Else standard delivery. Demo lookup, not live dark stores. No coupon.",
            WARM,
        ),
        (
            "5. COULD-HAVE — VIRTUAL TRY-ON",
            "Wanted: upload a photo, see the dress on you. Not shipped — no free API was good enough. Next iteration if a usable model exists. Not drawn as live.",
            AMBER_LT,
        ),
        (
            "DELIBERATELY NOT BUILT",
            "No numeric confidence %. No personalised fit from order history. No support chatbot. Thin evidence → “Too few reviews” / Need more info — never guess.",
            COOL,
        ),
    ]
    for i, (title, body, fill) in enumerate(feats):
        x = 0.32 + (i % 4) * 3.22
        card(s, x, 4.46, 3.10, 2.58, fill=fill, line=LINE)
        tb(s, x + 0.10, 4.52, 2.90, 0.50, [(title, {"bold": True, "color": NAVY})])
        tb(s, x + 0.10, 5.02, 2.90, 1.90, [(body, {})])
    footer(s, 4)


def slide_5(prs):
    s = blank(prs)
    header(s, "WHO WE TALKED TO AND WHY", "Three intenders, one freeze: they leave the app to decide")
    tb(
        s,
        0.32,
        0.94,
        12.7,
        0.42,
        [
            (
                "Three interview cards are directional composites from the dormant-item form themes and the catalog-tied pilot — not a statistically powered sample. Screener gate: ≥1 unsold wishlist item, not 5+.",
                {"italic": True, "color": MUTED},
            )
        ],
    )
    personas = [
        (
            "Aadyant, 26",
            "Occasion ethnic  ·  intent “Probably”",
            [
                ("Why wishlisted: liked the look; saved for a family function.", {}),
                ("Pain: listing photos vs threadwork/fabric in daylight.", {}),
                ("JTBD: know if it will look cheap in person before buying.", {}),
                ("Workaround: Google / cousin.", {}),
                ("“I still don’t know if the embroidery will look cheap in daylight.”", {"italic": True, "color": NAVY, "bold": True}),
            ],
            WARM,
        ),
        (
            "Meera, 28",
            "Two saved looks  ·  comparison freeze",
            [
                ("Why: comparing options already on the wishlist.", {}),
                ("Pain: gown vs kurta-set for the same evening; opens other apps.", {}),
                ("JTBD: pick one SKU without leaving Myntra.", {}),
                ("Workaround: YouTube unboxing + other saved piece.", {}),
                ("“I keep opening the other saved piece instead of deciding.”", {"italic": True, "color": NAVY, "bold": True}),
            ],
            COOL,
        ),
        (
            "Kabir, 24",
            "Fit freeze  ·  sitting 2+ weeks",
            [
                ("Why: liked the look; not sure on size.", {}),
                ("Pain: usual M vs reviews that say size up.", {}),
                ("JTBD: a clear size call so the item leaves the list.", {}),
                ("Workaround: order, then exchange.", {}),
                ("“I’ll order when I know whether to take M or L.”", {"italic": True, "color": NAVY, "bold": True}),
            ],
            WARM,
        ),
    ]
    for i, (name, ctx, lines, fill) in enumerate(personas):
        x = 0.32 + i * 4.30
        card(s, x, 1.40, 4.16, 3.55, fill=fill, line=LINE)
        rect(s, x, 1.40, 0.10, 3.55, NAVY)
        tb(s, x + 0.22, 1.46, 3.82, 0.32, [(name, {"bold": True, "color": NAVY})])
        tb(s, x + 0.22, 1.76, 3.82, 0.32, [(ctx, {"color": AMBER, "bold": True})])
        tb(s, x + 0.22, 2.12, 3.82, 2.70, lines)
    # Method row
    methods = [
        (
            "FORM  ·  PRIMARY",
            [
                [
                    ("Live Google Form  ·  ", {}),
                    ("open the viewform", {"url": FORM, "bold": True}),
                ],
                ("Recruited pilot, n = 6, monthly+ fashion, ≥1 unsold item. Why saved, still intend THIS item, main blocker now, unstick in 30 days, info needed, alternatives, outside the app, would an assistant help decide faster.", {}),
            ],
        ),
        (
            "1:1 INTERVIEW INSIGHTS",
            [
                (
                    "Confidence (quality / compare / fit) dominates these intenders. They leave the app to decide. Aadyant = quality, Meera = compare, Kabir = fit. Buying a different wishlisted SKU in 30 days is context only — not a gate.",
                    {},
                )
            ],
        ),
        (
            "SECONDARY  ·  USER BEHAVIOUR",
            [
                (
                    "Corpus: leave-app 10.8% of tagged rows; when they leave, other shopping apps are 60.9% of those mentions. Do not merge with 6/6. Reviews under-index the silent freeze — the form probes it directly.",
                    {},
                )
            ],
        ),
    ]
    for i, (title, body) in enumerate(methods):
        x = 0.32 + i * 4.30
        card(s, x, 5.05, 4.16, 2.02, fill=WHITE, line=NAVY)
        tb(s, x + 0.12, 5.10, 3.92, 0.32, [(title, {"bold": True, "color": AMBER})])
        tb(s, x + 0.12, 5.42, 3.92, 1.55, body)
    footer(s, 5)


def slide_6(prs):
    s = blank(prs)
    header(s, "WHAT PRIMARY RESEARCH CONFIRMED, AND CHANGED", "Corpus ranks price #2; these intenders are blocked by confidence")
    tb(
        s,
        0.32,
        0.94,
        12.7,
        0.32,
        [
            (
                "Confidence dominates for real intenders, even though price ranks #2 corpus-wide. Personas on slide 5 illustrate the cluster.",
                {"color": MUTED},
            )
        ],
    )
    card(s, 0.32, 1.32, 6.35, 4.55, fill=COOL, line=NAVY)
    tb(s, 0.48, 1.40, 6.05, 0.35, [("CONFIRMED", {"bold": True, "color": NAVY})])
    confirmed = [
        "6/6 left the app to do something before deciding — the confidence gap is real, not theoretical.",
        "6/6 main blockers sit in the confidence cluster: comparison, quality, occasion, fit. Price was not the main blocker for anyone in this pilot (0/6).",
        "4/6 named similar-buyer reviews/photos as what would unstick them; another 1/6 wanted clearer fit/sizing.",
        "4/6 said an in-app assistant would have helped them decide faster (2/6 Maybe, 0 No) — that is what the MVP delivers.",
        "Aadyant (quality), Meera (compare), Kabir (fit) are the three faces of the same freeze.",
        "Intent when saved: 4/6 Probably, 2/6 Yes, definitely — stuck in ambiguity, not abandoned.",
    ]
    tb(s, 0.48, 1.80, 6.05, 3.90, [(f"•  {c}", {}) for c in confirmed])
    card(s, 6.80, 1.32, 6.22, 4.55, fill=WARM, line=AMBER)
    tb(s, 6.96, 1.40, 5.90, 0.35, [("CHANGED", {"bold": True, "color": AMBER})])
    changed = [
        "Discovery corpus (n = 424): price is the #2 blocker overall at 23.8%.",
        "These intenders (n = 6): price is 0/6 as the main blocker; the confidence cluster is 6/6.",
        "So the target is the confidence cluster, not the price waiters — which is out of scope regardless, per the no-monetary-incentives constraint.",
        "Public reviews under-index freeze; the form is the second layer, not a combined rate.",
        "Read directionally, not statistically: n = 6 is a recruited pilot, not a powered sample.",
    ]
    tb(s, 6.96, 1.80, 5.90, 3.90, [(f"•  {c}", {}) for c in changed])
    card(s, 0.32, 6.00, 12.70, 1.08, fill=WHITE, line=LINE)
    tb(
        s,
        0.48,
        6.08,
        12.40,
        0.92,
        [
            (
                "Two evidence layers, never merged: public reviews mention leaving the app in 10.8% of tagged rows; survey respondents (n = 6) almost all did. Directional context, not one combined rate.",
                {},
            )
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    footer(s, 6)


def slide_7(prs):
    s = blank(prs)
    header(s, "THE PROBLEM, PRECISELY DEFINED", "No in-app way to resolve fit, quality, or compare on a saved item")
    cells = [
        ("SEGMENT", "24–30, Tier 1–2, monthly+ fashion. Must: ≥1 unsold saved item (not 5+) with “Probably” intent on ethnic/occasion or fashion-forward. 5+ / dormant 2+ weeks filters quotes, not the population."),
        ("PRODUCT OUTCOME", "Decision-confidence rate — funnel stage 3 (confidence resolution). Lift here is how WPCR moves."),
        ("ROOT CAUSE", "No in-app mechanism at reconsideration to resolve fit, quality-vs-photo, or comparison doubt. Users leave; the loop rarely closes."),
        ("WORKAROUNDS TODAY", "Order two sizes and return · ask friends · visit stores · search off-app reviews · let the wishlist grow forever."),
        ("USER VALUE", "Less anxiety, fewer returns, a faster yes-or-no on high-stakes occasion buys."),
        ("BUSINESS VALUE", "Higher-AOV ethnic demand already saved converts — without discounting."),
    ]
    for i, (k, v) in enumerate(cells):
        x = 0.32 + (i % 3) * 4.30
        y = 0.98 + (i // 3) * 1.38
        card(s, x, y, 4.16, 1.30, fill=WHITE, line=LINE)
        tb(s, x + 0.10, y + 0.04, 3.96, 0.26, [(k, {"bold": True, "color": AMBER})])
        tb(s, x + 0.10, y + 0.30, 3.96, 0.94, [(v, {})])
    tb(
        s,
        0.32,
        3.78,
        12.7,
        0.32,
        [
            (
                "Why this segment: already expressed intent (saved); high cost of being wrong (occasion, returns); in-app at reconsideration; the brief forbids discounts so we cannot buy conversion from price-waiters.",
                {"bold": True, "color": NAVY},
            )
        ],
    )
    tb(s, 0.32, 4.10, 12.7, 0.28, [("HYPOTHESES", {"bold": True, "color": AMBER})])
    hyps = [
        "1. Review-grounded verdicts on the card beat leaving for Google / YouTube.",
        "2. Compare 2–3 saved items beats open-ended browsing.",
        "3. A 30-min prompt (demo) can add urgency without a coupon.  Riskiest: they still will not buy without a physical / photo try-on — that is the named future, not this MVP.",
    ]
    tb(s, 0.32, 4.36, 12.7, 0.70, [(h, {}) for h in hyps])
    rows = [
        ["Player / behaviour", "What they do", "Gap this MVP fills"],
        ["Status quo", "Two sizes + return, friends, YouTube, Instagram", "No in-app close on the saved SKU"],
        [
            ("Amazon / Flipkart", {"url": AMAZON, "bold": True}),
            ("Size charts, Q&A", {"url": FLIPKART}),
            "Not Myntra wishlist; not occasion/threadwork grounded in this listing's reviews",
        ],
        [
            ("AJIO / Nykaa Fashion", {"url": NYKAA, "bold": True}),
            ("Same wishlist freeze", {"url": AJIO}),
            "Same confidence hole — no in-app resolve on the saved SKU",
        ],
        ["Virtual try-on (paid / studio)", "Photo on body", "Better for drape; not free here — listed as could-have"],
        [
            ("This MVP", {"bold": True, "fill": WARM, "color": NAVY}),
            ("In-wishlist agent + compare + demo 30-min", {"bold": True, "fill": WARM, "color": NAVY}),
            ("Grounded yes/no without money; try-on not shipped", {"bold": True, "fill": WARM, "color": NAVY}),
        ],
    ]
    table(s, 0.32, 5.08, 12.70, 2.00, rows, [3.2, 4.3, 5.2])
    footer(s, 7)


def slide_8(prs):
    s = blank(prs)
    header(s, "THE MVP: HOW IT SOLVES THIS", "The MVP closes that gap on the card — without a coupon")
    tb(
        s,
        0.32,
        0.94,
        12.7,
        0.30,
        [
            [
                ("Problem on slide 7 → solution on the wishlist card. Not a second product tour.  ", {"color": MUTED}),
                ("Deployed MVP", {"url": APP, "bold": True}),
            ]
        ],
    )
    rows = [
        ["Slide 7 cell", "MVP response"],
        ["Root cause: no in-app resolve", "Still deciding sits on the wishlist card"],
        ["Quality vs photo", "Reviews for this SKU; proof bullets beside the verdict"],
        ["Fit", "Usual size + Size up, then buy / Need more info"],
        ["Comparison", "Compare 2–3 saved items (hard cap)"],
        ["Workaround: leave app", "Keep the loop on Myntra"],
        ["Workaround: two sizes", "Verdict instead of guessing; returns as guardrail"],
        ["No urgency", "Express demo (secondary, not the hero)"],
        ["Physical try-on still missing", "Honest could-have; not faked"],
    ]
    table(s, 0.32, 1.28, 7.55, 4.55, rows, [3.3, 4.25])
    card(s, 8.02, 1.28, 5.00, 2.55, fill=WARM, line=AMBER)
    tb(s, 8.14, 1.34, 4.76, 0.30, [("WORKED EXAMPLE", {"bold": True, "color": AMBER})])
    tb(
        s,
        8.14,
        1.66,
        4.76,
        2.05,
        [
            [("Verdict: ", {"bold": True}), ("Size up, then buy", {"bold": True, "color": NAVY})],
            ("“Go one size up from your usual — this runs a little small.”", {"italic": True}),
            ("From reviews: several reviewers say it runs small and recommend ordering one size up.", {}),
            ("Then: Add to bag  or  Still not sure.", {}),
        ],
    )
    card(s, 8.02, 3.96, 5.00, 1.87, fill=COOL, line=NAVY)
    tb(s, 8.14, 4.02, 4.76, 0.30, [("DELIBERATELY NOT BUILT", {"bold": True, "color": NAVY})])
    tb(
        s,
        8.14,
        4.34,
        4.76,
        1.40,
        [
            ("No numeric confidence score of any kind.", {}),
            ("No personalised fit from the shopper's own order history.", {}),
            ("No customer-support chatbot.", {}),
            ("Thin evidence → Too few reviews / Need more info — never guesses.", {}),
        ],
    )
    card(s, 0.32, 5.96, 12.70, 1.12, fill=WHITE, line=LINE)
    tb(
        s,
        0.48,
        6.04,
        12.40,
        0.96,
        [
            (
                "Golden thread: WPCR → stage-3 decision confidence → discovery ranks quality/fit in and price out → research confirms confidence for intenders → in-wishlist assistant (+ compare, + demo 30-min). Express is a secondary urgency nudge. Try-on is named, not shipped.",
                {},
            )
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    footer(s, 8)


def slide_9(prs):
    s = blank(prs)
    header(s, "HOW WE’LL KNOW IT’S WORKING", "WPCR is the outcome; engagement, resolution, and guardrails tell us it’s real")
    tb(
        s,
        0.32,
        0.94,
        12.7,
        0.30,
        [
            (
                "Each metric has a definition, what it measures, and why it was chosen. No invented target percentages.",
                {"color": MUTED},
            )
        ],
    )
    rows = [
        ["Type", "Metric", "Definition", "What it measures", "Why we chose it"],
        [
            "North star",
            "WPCR",
            "% who buy ≥1 wishlisted item within 30 days of the add",
            "The brief’s business outcome",
            "Every slide exists to move this, without money",
        ],
        [
            "Leading",
            "Assistant engagement",
            "% of in-scope wishlist sessions that open Still deciding, Compare, or Express check",
            "Whether the tool is found at hesitation",
            "If they never open it, WPCR cannot move via this lever",
        ],
        [
            "Leading",
            "Resolution rate",
            "% of uses ending in add-to-bag, “this helped”, or a clear no",
            "Doubt actually closed",
            "A confident no is a success (better than freeze)",
        ],
        [
            "Leading",
            "Express attach (eligible pins only)",
            "% of eligible pin checks that bag as Express",
            "Whether urgency works without a coupon",
            "Optional; must not outrank confidence",
        ],
        [
            "Guardrail",
            "Return rate (fit/quality SKUs)",
            "Returns ÷ orders on items that used the assistant",
            "False confidence",
            "A wrong “Go for it” would show up here",
        ],
        [
            "Guardrail",
            "Wishlist-add rate",
            "Adds to wishlist per session / user",
            "Saving behaviour",
            "Tool should resolve doubt, not scare people off saving",
        ],
        [
            "Guardrail",
            "Fit / quality support contacts",
            "Contacts tagged fit or fabric/photo",
            "Whether in-app answers absorb those questions",
            "Expected to fall if the assistant works",
        ],
    ]
    table(s, 0.22, 1.26, 12.90, 5.80, rows, [1.35, 2.15, 3.55, 2.55, 3.30])
    footer(s, 9)


def slide_10(prs):
    s = blank(prs)
    header(s, "RISKS AND HOW WE’D DE-RISK THEM", "Solution-specific risks, and the steps already shipped to de-risk them")
    tb(
        s,
        0.32,
        0.94,
        12.7,
        0.28,
        [
            (
                "Most guardrails are already shipped; discoverability and try-on are planned. Standing constraint: never a discount, coupon, or cashback tool.",
                {"color": MUTED},
            )
        ],
    )
    rows = [
        ["Risk", "Why it matters", "Mitigation steps", "Status"],
        [
            "Hallucinated fit / quality",
            "Wrong buy → returns (guardrail)",
            "Writer sees only supplied reviews; verify figures vs aggregate; discard unsupported counts; fallback shows closest reviews, never a fake verdict",
            "SHIPPED",
        ],
        [
            "Thin evidence on new SKUs",
            "Guessing is worse than freeze",
            "Badge Too few reviews; verdict Need more info; demo includes a thin-review item",
            "SHIPPED",
        ],
        [
            "Distrust of “AI”",
            "They will not use it",
            "Proof bullets from real reviews beside every answer",
            "SHIPPED",
        ],
        [
            "Comparison overload",
            "More freeze",
            "Hard cap 2–3 items",
            "SHIPPED",
        ],
        [
            "Buried in a hub",
            "Engagement leading stays ~0",
            "Embed on the wishlist card at hesitation, not a separate hub",
            "PLANNED",
        ],
        [
            "30-min read as live ops",
            "Credibility hit",
            "Copy + data = demo hubs/prefixes; ineligible pin or missing size → standard delivery",
            "SHIPPED (demo)",
        ],
        [
            "Need physical try-on",
            "Riskiest hypothesis",
            "Do not fake try-on; name it as next if a free/good model exists",
            "PLANNED",
        ],
    ]
    # Color status column via per-cell opts
    styled = [rows[0]]
    for r in rows[1:]:
        status = r[3]
        fill = COOL if "SHIPPED" in status else WARM
        color = NAVY if "SHIPPED" in status else AMBER
        styled.append(
            [
                r[0],
                r[1],
                r[2],
                (status, {"bold": True, "color": color, "fill": fill, "align": PP_ALIGN.CENTER}),
            ]
        )
    table(s, 0.22, 1.24, 12.90, 5.82, styled, [2.35, 2.45, 6.00, 2.10])
    footer(s, 10)


def audit(path: Path):
    prs = Presentation(str(path))
    sizes = set()
    names = []
    urls = []
    texts = []
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    t = "".join(r.text for r in p.runs)
                    if t.strip():
                        texts.append((i, t))
                    for r in p.runs:
                        if r.font.size:
                            sizes.add(round(r.font.size.pt, 1))
                        if r.hyperlink and r.hyperlink.address:
                            urls.append((i, r.hyperlink.address))
            if sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            t = "".join(r.text for r in p.runs)
                            if t.strip():
                                texts.append((i, t))
                            for r in p.runs:
                                if r.font.size:
                                    sizes.add(round(r.font.size.pt, 1))
                                if r.hyperlink and r.hyperlink.address:
                                    urls.append((i, r.hyperlink.address))
    blob = "\n".join(t for _, t in texts)
    print(f"slides: {len(prs.slides)}")
    print(f"font sizes: {sorted(sizes)}")
    print(f"hyperlinks ({len(urls)}):")
    for i, u in urls:
        print(f"  slide {i}: {u}")
    for needle in ("n = 10", "n=10", "7/10", "10/10", "3/10", "5+ wishlist"):
        if needle.lower() in blob.lower():
            print(f"WARN leftover: {needle}")
    if "n = 6" not in blob and "n=6" not in blob:
        print("WARN: n=6 not found")
    old_form = "1Cie5So0fzPwsrBEB230FvdVQPjqDDbiQST6mWEK-ukE"
    if old_form in blob or any(old_form in u for _, u in urls):
        print("WARN: old converter form still linked")
    print(f"wrote {path}")


def main():
    make_blocker_chart(CHART)
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)
    slide_6(prs)
    slide_7(prs)
    slide_8(prs)
    slide_9(prs)
    slide_10(prs)
    prs.save(str(OUT))
    audit(OUT)


if __name__ == "__main__":
    main()
