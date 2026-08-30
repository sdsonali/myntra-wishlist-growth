"""Dump a PPTX structure to a UTF-8 text file for inspection."""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

SRC = Path(sys.argv[1] if len(sys.argv) > 1 else r"C:\Users\itch\Downloads\Wishlist_Confidence_Gap_Deck.pptx")
OUT = Path(__file__).resolve().parent / "deck_dump.txt"


def color_of(run):
    try:
        return str(run.font.color.rgb)
    except Exception:
        return "-"


def main():
    prs = Presentation(str(SRC))
    lines = [
        f"file: {SRC}",
        f"slides: {len(prs.slides)}  size: {Emu(prs.slide_width).inches:.2f}x{Emu(prs.slide_height).inches:.2f} in",
        "",
    ]
    sizes = set()
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"===== SLIDE {i} =====")
        for sh in slide.shapes:
            lines.append(
                f"  [{sh.shape_id}] {sh.name} {sh.shape_type} "
                f"L={Emu(sh.left).inches:.2f} T={Emu(sh.top).inches:.2f} "
                f"W={Emu(sh.width).inches:.2f} H={Emu(sh.height).inches:.2f}"
            )
            try:
                fill = sh.fill
                if fill.type is not None and fill.type == 1:
                    lines.append(f"        fill={fill.fore_color.rgb}")
            except Exception:
                pass
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    txt = "".join(r.text for r in p.runs)
                    if not txt.strip():
                        continue
                    for r in p.runs:
                        if r.font.size:
                            sizes.add(r.font.size.pt)
                    meta = " | ".join(
                        f"{(r.font.size.pt if r.font.size else '-')}pt "
                        f"{'B' if r.font.bold else ''} {color_of(r)}"
                        for r in p.runs
                    )
                    lines.append(f"        T: {txt}")
                    lines.append(f"           ({meta})")
            if sh.has_table:
                tbl = sh.table
                lines.append(f"        TABLE {len(tbl.rows)}x{len(tbl.columns)}")
                for row in tbl.rows:
                    cells = [c.text.replace("\n", " / ") for c in row.cells]
                    lines.append("          | " + " | ".join(cells))
        lines.append("")
    lines.append(f"ALL FONT SIZES USED: {sorted(sizes)}")
    OUT.write_text("\n".join(lines), encoding="utf-8")
    print(f"wrote {OUT}")
    print("font sizes:", sorted(sizes))


if __name__ == "__main__":
    main()
