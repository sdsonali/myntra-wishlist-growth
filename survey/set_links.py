"""Replace the deck's placeholder artefact links with real URLs and re-export the PDF.

Usage (any subset of the three flags):

    pm-env\\Scripts\\python.exe survey\\set_links.py ^
        --discovery https://your-app.streamlit.app/ ^
        --mvp https://your-app.streamlit.app/ ^
        --survey https://docs.google.com/forms/d/e/XXXX/viewform

Rewrites both the hyperlink target and the visible "[LINK: ...]" placeholder text,
then exports a fresh PDF next to the .pptx.
"""

import argparse
import sys
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
PPTX = ROOT / "deck" / "Wishlist_Confidence_Gap_Deck_v2.pptx"
PDF = ROOT / "deck" / "Myntra_Wishlist_Confidence_Gap_FINAL.pdf"

# placeholder URL fragment -> (cli arg name, visible label)
ARTEFACTS = {
    "discovery-tab1": ("discovery", "AI Discovery Engine"),
    "mvp-tab2": ("mvp", "Deployed MVP"),
    "survey-form": ("survey", "Pilot survey form"),
}


def classify(text, address):
    """Return the artefact key implied by a run's hyperlink target or its text."""
    for fragment, _ in ARTEFACTS.items():
        if address and fragment in address:
            return fragment
    lowered = (text or "").lower()
    if "discovery" in lowered:
        return "discovery-tab1"
    if "mvp" in lowered:
        return "mvp-tab2"
    if "form" in lowered or "survey" in lowered:
        return "survey-form"
    return None


def iter_text_frames(shapes):
    for shape in shapes:
        if shape.shape_type == 6:  # group
            yield from iter_text_frames(shape.shapes)
            continue
        if shape.has_text_frame:
            yield shape.text_frame
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--discovery")
    ap.add_argument("--mvp")
    ap.add_argument("--survey")
    ap.add_argument("--no-pdf", action="store_true", help="skip the PDF export")
    args = ap.parse_args()

    urls = {
        key: getattr(args, arg)
        for key, (arg, _) in ARTEFACTS.items()
        if getattr(args, arg)
    }
    if not urls:
        ap.error("give at least one of --discovery / --mvp / --survey")

    prs = Presentation(PPTX)
    changed = []

    for index, slide in enumerate(prs.slides, 1):
        for frame in iter_text_frames(slide.shapes):
            for para in frame.paragraphs:
                for position, run in enumerate(para.runs):
                    address = run.hyperlink.address
                    is_placeholder = "[LINK:" in (run.text or "")
                    if not is_placeholder and not (address and "example.com" in address):
                        continue

                    key = classify(run.text, address)
                    if key not in urls:
                        continue

                    if is_placeholder:
                        label = ARTEFACTS[key][1]
                        # slides that already print "Deployed MVP:" would otherwise read twice
                        preceding = "".join(r.text or "" for r in para.runs[:position]).lower()
                        run.text = "Open link" if label.lower() in preceding else label
                    run.hyperlink.address = urls[key]
                    changed.append((index, key, run.text))

    if not changed:
        print("nothing to change — no placeholders matched")
        return

    prs.save(PPTX)
    for slide_no, key, label in changed:
        print(f"slide {slide_no}: {key} -> {urls[key]}  (text: {label!r})")
    print(f"\nsaved {PPTX}")

    if args.no_pdf:
        return

    try:
        import win32com.client  # noqa: F401  (import kept local; COM is optional)
    except ImportError:
        print("\npywin32 not installed — export the PDF from PowerPoint manually")
        return

    import win32com.client as com

    app = com.Dispatch("PowerPoint.Application")
    pres = app.Presentations.Open(str(PPTX), True, False, False)
    if PDF.exists():
        PDF.unlink()
    pres.SaveCopyAs(str(PDF), 32)
    pres.Close()
    app.Quit()
    print(f"exported {PDF}")


if __name__ == "__main__":
    sys.exit(main())
