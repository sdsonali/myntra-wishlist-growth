"""
NextLeap 10-slide deck (14pt everywhere, no fellow name).

Fill LINKS below, then:
  python survey/generate_deck.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Public artifacts — paste viewable URLs (test in incognito).
LINKS = {
    "discovery": "",  # Streamlit Tab 1
    "mvp": "",  # Streamlit Tab 2 (same app is fine)
    "survey": "",  # Google Form viewform URL
}

OUT = Path(__file__).resolve().parent.parent / "deck" / "Myntra_Wishlist_Confidence_Gap.pptx"

NAVY = RGBColor(0x1E, 0x3A, 0x5F)
INK = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x4A, 0x4A, 0x4A)
BLUE = RGBColor(0x1F, 0x4E, 0x79)
AMBER = RGBColor(0xB8, 0x5C, 0x38)
BG = RGBColor(0xF6, 0xF4, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC8, 0xC4, 0xBC)
LINK = RGBColor(0x1F, 0x4E, 0x79)

W = Inches(13.333)
H = Inches(7.5)
FS = Pt(14)
FONT = "Calibri"


def _set_run(run, text, *, bold=False, color=INK, italic=False, underline=False):
    run.text = text
    run.font.name = FONT
    run.font.size = FS
    run.font.bold = bold
    run.font.italic = italic
    run.font.underline = underline
    run.font.color.rgb = color


def _fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def _box(slide, l, t, w, h, fill=None, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.adjustments[0] = 0.08
    if fill:
        _fill(sh, fill)
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    elif fill:
        sh.line.fill.background()
    return sh


def _rect(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    _fill(sh, fill)
    return sh


def _tf(shape, *, anchor=MSO_ANCHOR.TOP):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    return tf


def add_text(slide, l, t, w, h, paragraphs, *, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    """paragraphs: list of str or (str, dict)."""
    sh = slide.shapes.add_textbox(l, t, w, h)
    tf = _tf(sh, anchor=anchor)
    first = True
    for item in paragraphs:
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = opts.get("align", align)
        p.space_after = Pt(opts.get("space_after", 4))
        p.level = 0
        run = p.add_run()
        _set_run(
            run,
            text,
            bold=opts.get("bold", False),
            color=opts.get("color", INK),
            italic=opts.get("italic", False),
            underline=opts.get("underline", False),
        )
        url = opts.get("url")
        if url:
            run.hyperlink.address = url
            _set_run(run, text, bold=opts.get("bold", False), color=LINK, underline=True)
    return sh


def add_mixed(slide, l, t, w, h, runs, *, space_after=4):
    """One paragraph with mixed runs: list of (text, opts)."""
    sh = slide.shapes.add_textbox(l, t, w, h)
    tf = _tf(sh)
    p = tf.paragraphs[0]
    p.space_after = Pt(space_after)
    for text, opts in runs:
        run = p.add_run()
        url = opts.get("url")
        _set_run(
            run,
            text,
            bold=opts.get("bold", False),
            color=LINK if url else opts.get("color", INK),
            italic=opts.get("italic", False),
            underline=bool(url) or opts.get("underline", False),
        )
        if url:
            run.hyperlink.address = url
    return sh


def link_or_note(key, label):
    url = (LINKS.get(key) or "").strip()
    if url:
        return (label, {"url": url, "bold": True})
    return (f"{label} (add public URL in generate_deck.py LINKS)", {"italic": True, "color": MUTED})


def bg(slide):
    _rect(slide, 0, 0, W, H, BG)
    _rect(slide, 0, 0, Inches(0.12), H, NAVY)


def title_bar(slide, title, n):
    add_text(
        slide,
        Inches(0.4),
        Inches(0.18),
        Inches(11.6),
        Inches(0.85),
        [(title, {"bold": True, "color": NAVY})],
    )
    add_text(
        slide,
        Inches(12.05),
        Inches(0.22),
        Inches(1.0),
        Inches(0.4),
        [(f"{n} / 10", {"color": MUTED, "align": PP_ALIGN.RIGHT})],
        align=PP_ALIGN.RIGHT,
    )
    _rect(slide, Inches(0.4), Inches(1.02), Inches(12.5), Pt(1.5), LINE)


def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    return slide


def slide_1(prs):
    s = blank(prs)
    _rect(s, 0, 0, W, H, NAVY)
    _rect(s, 0, Inches(5.85), W, Inches(1.65), AMBER)
    add_text(
        s,
        Inches(0.7),
        Inches(2.0),
        Inches(12),
        Inches(1.2),
        [("Closing the Wishlist Confidence Gap", {"bold": True, "color": WHITE})],
    )
    add_text(
        s,
        Inches(0.7),
        Inches(3.2),
        Inches(12),
        Inches(1.4),
        [
            (
                "Improving Wishlist-to-Purchase Conversion on Myntra — without discounts",
                {"color": WHITE},
            )
        ],
    )
    add_text(
        s,
        Inches(0.7),
        Inches(6.2),
        Inches(12),
        Inches(0.8),
        [("Product Manager, Growth Team", {"bold": True, "color": WHITE})],
    )


def slide_2(prs):
    s = blank(prs)
    title_bar(
        s,
        "Re-engagement is not the bottleneck — decision confidence is the lever that moves WPCR",
        2,
    )
    add_text(
        s,
        Inches(0.4),
        Inches(1.12),
        Inches(12.5),
        Inches(0.4),
        [
            (
                "North star (WPCR): % of users who purchase ≥1 wishlisted item within 30 days of adding it.",
                {"color": MUTED},
            )
        ],
    )
    stages = [
        ("1. Re-engagement", "% of wishlisted items revisited in 30 days", False),
        ("2. Decision confidence", "% of revisited items where fit / quality / comparison doubt is resolved", True),
        ("3. Cart conversion", '% of "resolved" items added to cart', False),
        ("4. Checkout", "% of cart items purchased", False),
    ]
    x0 = Inches(0.4)
    gap = Inches(0.18)
    bw = Inches(3.0)
    for i, (name, desc, hot) in enumerate(stages):
        x = x0 + i * (bw + gap)
        box = _box(s, x, Inches(1.6), bw, Inches(1.85), fill=AMBER if hot else WHITE, line=NAVY if hot else LINE)
        add_text(
            s,
            x,
            Inches(1.68),
            bw,
            Inches(1.7),
            [
                (name, {"bold": True, "color": WHITE if hot else NAVY}),
                (desc, {"color": WHITE if hot else MUTED}),
                ("PRIMARY LEVER" if hot else " ", {"bold": True, "color": WHITE if hot else WHITE}),
            ],
        )
    add_text(
        s,
        Inches(0.4),
        Inches(3.6),
        Inches(12.5),
        Inches(0.35),
        [("Segmentation lens (labels, not colour alone)", {"bold": True, "color": NAVY})],
    )
    segs = [
        "Bookmarkers (mood-board) vs Genuine Intenders",
        "Ethnic / occasion wear (fit-sensitive, high AOV) vs basics (fit-agnostic, low AOV)",
        "Price-waiters (out of scope — no monetary levers) vs confidence-blocked (in scope)",
    ]
    add_text(s, Inches(0.4), Inches(3.95), Inches(12.5), Inches(1.4), [(f"•  {t}", {}) for t in segs])
    call = _box(s, Inches(0.4), Inches(5.45), Inches(12.5), Inches(1.55), fill=WHITE, line=AMBER)
    add_text(
        s,
        Inches(0.5),
        Inches(5.55),
        Inches(12.3),
        Inches(1.35),
        [
            (
                "Callout: target shoppers already visit 2–4×/month, so stage 1 is not the gap. "
                "Wishlisted demand dies at stage 2 — doubt is not resolved in-app, so cart and checkout never get a chance.",
                {"color": INK},
            )
        ],
    )


def slide_3(prs):
    s = blank(prs)
    title_bar(
        s,
        "AI discovery: quality is the loudest complaint; fit is real; price exists but is out of scope",
        3,
    )
    add_text(
        s,
        Inches(0.4),
        Inches(1.12),
        Inches(12.5),
        Inches(0.55),
        [
            (
                "Workflow: Play Store + App Store + YouTube (n=550 tagged reviews) → LLM theme clustering → tagging → mention share. Reviews complain about received products; they under-index silent wishlist freeze — primary research fills that gap.",
                {"color": MUTED},
            )
        ],
    )
    steps = ["Sources", "Cluster", "Tag", "Quantify"]
    for i, lab in enumerate(steps):
        x = Inches(0.4) + i * Inches(2.05)
        _box(s, x, Inches(1.75), Inches(1.85), Inches(0.55), fill=NAVY)
        add_text(s, x, Inches(1.8), Inches(1.85), Inches(0.45), [(lab, {"bold": True, "color": WHITE, "align": PP_ALIGN.CENTER})], align=PP_ALIGN.CENTER)
        if i < 3:
            add_text(s, x + Inches(1.75), Inches(1.8), Inches(0.35), Inches(0.45), [("→", {"bold": True, "color": AMBER})])
    rows = [
        ("Quality / authenticity doubt (incl. fabric, finish)", "47.8%", "IN SCOPE", True),
        ("Price change / deal-waiting", "13.9%", "OUT OF SCOPE", False),
        ("Fit / size uncertainty", "13.0%", "IN SCOPE", True),
        ("Found an alternative (comparison)", "5.2%", "IN SCOPE", True),
        ("Occasion mismatch", "2.6%", "IN SCOPE", True),
    ]
    y = Inches(2.5)
    add_text(
        s,
        Inches(0.4),
        y,
        Inches(8.4),
        Inches(0.35),
        [("Top opportunity areas (mention share)", {"bold": True, "color": NAVY})],
    )
    y = Inches(2.9)
    for label, pct, tag, hot in rows:
        _box(s, Inches(0.4), y, Inches(8.5), Inches(0.52), fill=WHITE, line=LINE)
        add_text(s, Inches(0.5), y + Inches(0.05), Inches(5.4), Inches(0.42), [(label, {})])
        add_text(s, Inches(5.9), y + Inches(0.05), Inches(1.1), Inches(0.42), [(pct, {"bold": True, "color": NAVY})])
        add_text(
            s,
            Inches(7.05),
            y + Inches(0.05),
            Inches(1.75),
            Inches(0.42),
            [(tag, {"bold": True, "color": AMBER if hot else MUTED})],
        )
        y += Inches(0.56)
    _box(s, Inches(9.15), Inches(2.9), Inches(3.75), Inches(3.7), fill=WHITE, line=NAVY)
    add_text(
        s,
        Inches(9.25),
        Inches(3.05),
        Inches(3.55),
        Inches(3.4),
        [
            ("How to read this", {"bold": True, "color": NAVY}),
            ("Quality is the loudest public complaint. Fit and comparison are smaller in review text but real. Price is present — we still cannot use monetary levers.", {}),
            ("Also tagged: comparison 11.8%, external research 10.5% of the corpus.", {}),
            link_or_note("discovery", "Open AI Discovery Engine"),
        ],
    )


def slide_4(prs):
    s = blank(prs)
    title_bar(
        s,
        'Pilot (n=10): “Probably” buyers leave the app to decide — no in-app way to resolve doubt',
        4,
    )
    add_text(
        s,
        Inches(0.4),
        Inches(1.12),
        Inches(12.5),
        Inches(0.7),
        [
            (
                "Pilot Google Form, n=10 (one named item each). Target we designed for: 24–30, Tier 1–2, 2–4×/month, 5+ wishlist, dormant 2+ weeks. "
                "Named items: 4/10 ethnic/occasion (Rakhi dress, wedding dress, kurta set, kurta); also sneakers, heels, shirts, cosmetics, top.",
                {"color": MUTED},
            )
        ],
    )
    left = [
        ("Main blocker (self-reported)", {"bold": True, "color": NAVY}),
        ("• Price / budget — 3/10 (out of scope)", {}),
        ("• Comparison paralysis — 2/10", {}),
        ("• Quality doubt — 2/10", {}),
        ("• Occasion / style fit — 2/10", {}),
        ("• Fit doubt — 1/10", {}),
        ("Combined confidence blockers 7/10 once price is set aside.", {"bold": True}),
    ]
    add_text(s, Inches(0.4), Inches(1.85), Inches(6.1), Inches(2.6), left)
    right = [
        ("Purchase intent", {"bold": True, "color": NAVY}),
        ('6/10 said “Probably” — stuck in ambiguity, neither committed nor abandoned. That is the target zone.', {}),
        ("Outside the app — 10/10 did at least one", {"bold": True, "color": NAVY}),
        ("Compare prices 6/10 · visit a store 6/10 · Instagram 3/10 · friends/family 3/10 · Google/YouTube 3/10. The app is a save, not a decision environment.", {}),
        ("What would unstick them: 7/10 wanted reviews/photos from similar buyers. Coping: order two and keep the fit; size-chart then exchange; leave it until payday.", {}),
    ]
    add_text(s, Inches(6.7), Inches(1.85), Inches(6.2), Inches(3.4), right)
    _box(s, Inches(0.4), Inches(5.5), Inches(12.5), Inches(1.5), fill=WHITE, line=AMBER)
    add_text(
        s,
        Inches(0.5),
        Inches(5.6),
        Inches(12.3),
        Inches(1.3),
        [
            (
                "Caveat: n=10 is directional, not powered. Screener was leaky — Q2/Q3 blank for 7/10; one respondent shops Rarely. Counts illustrate themes, not population rates. Research decision: pursue the 7/10 confidence cluster, not the 3/10 price-waiters.",
                {},
            ),
            link_or_note("survey", "Open screening survey (Google Form)"),
        ],
    )


def slide_5(prs):
    s = blank(prs)
    title_bar(
        s,
        "The gap is in-app confidence at reconsideration — not lack of desire to buy",
        5,
    )
    add_text(
        s,
        Inches(0.4),
        Inches(1.12),
        Inches(12.5),
        Inches(1.15),
        [
            (
                "Segment: 24–30, Tier 1–2, 2–4×/month, impulsive savers with 5+ items sitting 2+ weeks, “Probably” intent, especially ethnic/occasion and fashion-forward (not basics). "
                "Product outcome: Decision Confidence Rate (funnel stage 2) on those SKUs.",
                {},
            ),
            (
                "Root cause: after an impulsive save there is no in-app way, at revisit, to resolve fit, quality/authenticity (fabric, threadwork vs photo), comparison across similar saved options, or occasion-fit. Users leave; most never complete the loop.",
                {},
            ),
        ],
    )
    add_text(
        s,
        Inches(0.4),
        Inches(2.35),
        Inches(12.5),
        Inches(1.35),
        [
            (
                "Workarounds today: order multiple sizes and return extras; screenshot to friends; visit stores; read size-chart threads off-app.",
                {},
            ),
            (
                "User value: less anxiety, fewer return hassles, faster confident decisions on high-stakes occasion buys.",
                {},
            ),
            (
                "Business value: ethnic/occasion AOV is higher than basics, so a modest confidence lift moves more GMV from demand already on the wishlist — no acquisition, no discount. If answers are grounded, returns and fit tickets can fall (guardrail: false confidence must not raise returns).",
                {},
            ),
        ],
    )
    add_text(s, Inches(0.4), Inches(3.85), Inches(12.5), Inches(0.35), [("How the argument was built", {"bold": True, "color": NAVY})])
    flow = [
        "Business metric",
        "Product outcomes",
        "AI discovery",
        "Primary research",
        "Problem definition",
    ]
    notes = [
        "WPCR",
        "Stage 2",
        "Fit / quality / compare",
        "Ethnic + Probably + leave-app",
        "In-app confidence gap",
    ]
    for i, (a, b) in enumerate(zip(flow, notes)):
        x = Inches(0.4) + i * Inches(2.5)
        _box(s, x, Inches(4.25), Inches(2.3), Inches(1.55), fill=WHITE, line=NAVY)
        add_text(
            s,
            x,
            Inches(4.35),
            Inches(2.3),
            Inches(1.35),
            [(a, {"bold": True, "color": NAVY}), (b, {"color": MUTED})],
        )


def slide_6(prs):
    s = blank(prs)
    title_bar(
        s,
        "Ethnic-wear decision confidence is the lever that fits the data and the no-discount constraint",
        6,
    )
    reasons = [
        (
            "Largest combined blocker we are allowed to solve",
            "Pilot: confidence cluster 7/10 vs price 3/10. Discovery: quality 47.8% + fit 13% of blocker mentions. Ethnic/forward pieces are less “safe” than basics (fit + threadwork/fabric vs photo).",
        ),
        (
            "No margin giveaway",
            "Confidence tools do not require discounts — they match the “no monetary incentives” brief.",
        ),
        (
            "Instrumentable stage",
            "Maps directly to Decision Confidence Rate on the WPCR funnel (slide 2), so we can measure the product outcome we claim to move.",
        ),
        (
            "Data already exists",
            "Reviews, return-reason text, and size-chart data can power the MVP. We do not need a new data-collection program to ship a first version.",
        ),
        (
            "Why ethnic/forward, not basics or price",
            "Basics are lower AOV and lower uncertainty — users already “just order.” Price-waiters are real (salary, sale) but forbidden by the brief. Bookmarking is not purchase intent.",
        ),
    ]
    y = Inches(1.2)
    for title, body in reasons:
        _box(s, Inches(0.4), y, Inches(12.5), Inches(1.05), fill=WHITE, line=LINE)
        add_text(
            s,
            Inches(0.55),
            y + Inches(0.08),
            Inches(12.2),
            Inches(0.9),
            [(title, {"bold": True, "color": NAVY}), (body, {})],
        )
        y += Inches(1.12)


def slide_7(prs):
    s = blank(prs)
    title_bar(
        s,
        "MVP: an AI layer inside the wishlist that answers fit, quality, and comparison without leaving the app",
        7,
    )
    add_text(
        s,
        Inches(0.4),
        Inches(1.12),
        Inches(12.5),
        Inches(0.4),
        [("Fit & Confidence Assistant — not a new destination; it sits on the existing wishlist.", {"color": MUTED})],
    )
    feats = [
        ("Confidence score", 'e.g. “87% of similar buyers say true to size,” from aggregated review / return-reason data.'),
        ("Ask-a-question", "Grounded answers on fit, fabric/threadwork vs photo, and occasion — replaces Google, Instagram, and WhatsApp-a-friend."),
        ("Personalized fit", "Cross-reference the user’s past order / fit feedback on similar cuts or brands when available."),
        ("Comparison assist", "If 2–3 similar items are saved, a tight side-by-side to break comparison paralysis."),
    ]
    y = Inches(1.55)
    for t, b in feats:
        add_text(
            s,
            Inches(0.4),
            y,
            Inches(7.3),
            Inches(0.85),
            [(t, {"bold": True, "color": NAVY}), (b, {})],
        )
        y += Inches(0.88)
    # Mock wishlist card
    _box(s, Inches(8.0), Inches(1.55), Inches(4.9), Inches(4.55), fill=WHITE, line=NAVY)
    add_text(s, Inches(8.15), Inches(1.65), Inches(4.6), Inches(0.4), [("Wishlist · mockup (prototype)", {"bold": True, "color": MUTED})])
    _box(s, Inches(8.25), Inches(2.15), Inches(4.4), Inches(1.15), fill=NAVY)
    add_text(
        s,
        Inches(8.35),
        Inches(2.25),
        Inches(4.2),
        Inches(0.95),
        [
            ("Kurta set · occasion wear", {"bold": True, "color": WHITE}),
            ("Saved 3 weeks ago · intent: Probably", {"color": WHITE}),
        ],
    )
    _box(s, Inches(8.25), Inches(3.45), Inches(4.4), Inches(0.85), fill=AMBER)
    add_text(
        s,
        Inches(8.35),
        Inches(3.52),
        Inches(4.2),
        Inches(0.75),
        [
            ("Confidence 87%  ·  true to size (labelled)", {"bold": True, "color": WHITE}),
        ],
    )
    add_text(
        s,
        Inches(8.25),
        Inches(4.4),
        Inches(4.4),
        Inches(1.5),
        [
            ("Ask: Will this work for a daytime wedding?", {}),
            ("Answer grounded in reviews + similar-buyer fit notes. Provenance shown, not a black-box claim.", {"color": MUTED}),
        ],
    )
    add_text(s, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.7), [link_or_note("mvp", "Open deployed Fit & Confidence Assistant (Tab 2)")])


def slide_8(prs):
    s = blank(prs)
    title_bar(
        s,
        "WPCR is the north star; we watch whether doubt is resolved — and that we do not inflate returns",
        8,
    )
    rows = [
        ("North star", "WPCR (30-day wishlist-to-purchase)", "Stage 4 outcome of the funnel (slide 2). Lagging, but the business result we owe."),
        ("Leading · adoption", "Fit & Confidence Assistant engagement", "% of wishlist visits that use the tool. If unused, we cannot move stage 2."),
        ("Leading · resolution", "Decision confidence rate", '% of interactions ending in add-to-cart or a “this helped” signal. Checks we resolve doubt, not just get clicks. Ties to slide 4 “Probably” zone.'),
        ("Guardrail", "Return rate must not increase", "False confidence would show up as more returns — the opposite of the user-value case."),
        ("Guardrail", "Wishlist add rate must not drop", "The tool must not punish browsing/saving. Bookmarkers can still bookmark."),
        ("Guardrail", "Fit-related support / contact rate", "Expected to fall if in-app answers replace tickets and WhatsApp workarounds (slide 4)."),
    ]
    y = Inches(1.15)
    for a, b, c in rows:
        _box(s, Inches(0.4), y, Inches(12.5), Inches(0.88), fill=WHITE, line=LINE)
        add_text(s, Inches(0.5), y + Inches(0.08), Inches(2.3), Inches(0.72), [(a, {"bold": True, "color": AMBER})])
        add_text(s, Inches(2.85), y + Inches(0.08), Inches(3.6), Inches(0.72), [(b, {"bold": True, "color": NAVY})])
        add_text(s, Inches(6.55), y + Inches(0.08), Inches(6.2), Inches(0.72), [(c, {})])
        y += Inches(0.95)


def slide_9(prs):
    s = blank(prs)
    title_bar(
        s,
        "We de-risk cold start, hallucination, trust, hiding, and comparison overload before scale",
        9,
    )
    rows = [
        (
            "Cold start on new / low-review ethnic SKUs",
            "Fallback to brand-level size charts + community Q&A. Say “insufficient item-level evidence” instead of inventing a score.",
        ),
        (
            "AI hallucination on fit / quality",
            "Ground strictly in review and return text. Show provenance. Use confidence caveats, never absolute claims.",
        ),
        (
            "Low trust in AI-only answers",
            "Surface real user reviews beside the summary — not instead of them. Label model output as a synthesis.",
        ),
        (
            "Low discoverability",
            "Embed at the point of hesitation inside the wishlist row, not a separate hub or after-checkout module.",
        ),
        (
            "Comparison increases fatigue",
            "Cap at the top 2–3 similar wishlisted items. No open-ended catalogue compare.",
        ),
    ]
    y = Inches(1.18)
    for risk, mit in rows:
        _box(s, Inches(0.4), y, Inches(12.5), Inches(1.05), fill=WHITE, line=LINE)
        add_text(s, Inches(0.55), y + Inches(0.08), Inches(4.4), Inches(0.9), [(risk, {"bold": True, "color": NAVY})])
        add_text(s, Inches(5.1), y + Inches(0.08), Inches(7.6), Inches(0.9), [(mit, {})])
        y += Inches(1.1)


def slide_10(prs):
    s = blank(prs)
    title_bar(
        s,
        "Resolve fit, quality, and comparison doubt at hesitation to lift ethnic-wear WPCR without touching price",
        10,
    )
    _box(s, Inches(0.4), Inches(1.25), Inches(12.5), Inches(2.2), fill=WHITE, line=NAVY)
    add_text(
        s,
        Inches(0.6),
        Inches(1.4),
        Inches(12.1),
        Inches(1.95),
        [
            ("One-line recap", {"bold": True, "color": NAVY}),
            (
                "By resolving fit, quality (including online vs real fabric/threadwork), and comparison at the moment of hesitation on high-AOV ethnic/fashion-forward wishlist items — validated by discovery (quality+fit) and the n=10 pilot (7/10 confidence, 6/10 Probably, 10/10 leave-app) — we lift decision-confidence rate and thus WPCR, without touching price.",
                {},
            ),
        ],
    )
    add_text(
        s,
        Inches(0.4),
        Inches(3.65),
        Inches(12.5),
        Inches(2.6),
        [
            ("Supporting artefacts (must be viewable logged out)", {"bold": True, "color": NAVY}),
            link_or_note("discovery", "AI Discovery Engine"),
            link_or_note("mvp", "Deployed MVP — Fit & Confidence Assistant"),
            link_or_note("survey", "Primary-research screening survey"),
            ("Export this deck to PDF for submission if the course asks for a PDF.", {"color": MUTED, "italic": True}),
        ],
    )
    add_text(
        s,
        Inches(0.4),
        Inches(6.4),
        Inches(12.5),
        Inches(0.55),
        [("Product Manager, Growth Team  ·  no monetary incentives", {"color": MUTED})],
    )


def main():
    prs = new_prs()
    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)
    slide_6(prs)
    slide_7(prs)
    slide_8(prs)
    slide_9(prs)
    slide_10(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
