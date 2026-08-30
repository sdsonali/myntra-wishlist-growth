"""
Rebuild slide 7 (MVP) of the Claude deck as a Myntra-app-style wishlist mock.

  python survey/enhance_deck.py                 # rebuild slide 7
  python survey/enhance_deck.py --font14        # also push slide copy to 14pt

Shapes named "mock_*" are UI chrome inside the phone and are left small by --font14.
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
SRC = Path(r"C:\Users\itch\Downloads\Wishlist_Confidence_Gap_Deck.pptx")
OUT = ROOT / "deck" / "Wishlist_Confidence_Gap_Deck_v2.pptx"
IMG_DIR = ROOT / "data" / "mvp_images"

NAVY = RGBColor(0x07, 0x29, 0x3E)
NAVY2 = RGBColor(0x0B, 0x3D, 0x5C)
AMBER = RGBColor(0xE8, 0xA3, 0x3D)
AMBER_DK = RGBColor(0xB9, 0x7A, 0x1E)
LIGHT = RGBColor(0xF4, 0xF7, 0xF9)
MUTED = RGBColor(0x5C, 0x6B, 0x73)
INK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xC7, 0xD6, 0xE0)
LINE = RGBColor(0xEA, 0xEA, 0xEC)
PINK = RGBColor(0xFF, 0x3F, 0x6C)
BEZEL = RGBColor(0x22, 0x22, 0x22)

FONT = "Calibri"


def clear(slide):
    tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        tree.remove(shape._element)


def rect(slide, l, t, w, h, fill=None, line=None, rounded=False, adj=0.06, name=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h),
    )
    if rounded:
        shape.adjustments[0] = adj
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    if name:
        shape.name = name
    return shape


def text(slide, l, t, w, h, runs, *, size=12.5, color=INK, bold=False,
         align=PP_ALIGN.LEFT, italic=False, space=2, name=None):
    """runs: str, or list of str (paragraphs), or list of (str, dict)."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = 0
    tf.margin_bottom = 0
    items = [runs] if isinstance(runs, str) else runs
    first = True
    for item in items:
        body, opts = (item, {}) if isinstance(item, str) else item
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = opts.get("align", align)
        para.space_after = Pt(opts.get("space", space))
        run = para.add_run()
        run.text = body
        font = run.font
        font.name = FONT
        font.size = Pt(opts.get("size", size))
        font.bold = opts.get("bold", bold)
        font.italic = opts.get("italic", italic)
        font.color.rgb = opts.get("color", color)
    if name:
        box.name = name
    return box


def text_mixed(slide, l, t, w, h, runs, *, size=11.0):
    """One paragraph, several runs: [(text, {bold, color}), ...]."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_top = 0
    tf.margin_bottom = 0
    para = tf.paragraphs[0]
    for body, opts in runs:
        run = para.add_run()
        run.text = body
        run.font.name = FONT
        run.font.size = Pt(opts.get("size", size))
        run.font.bold = opts.get("bold", False)
        run.font.color.rgb = opts.get("color", INK)
    return box


def pill(slide, l, t, w, h, label, *, fill, fg, size=8.0, bold=True, line=None, name="mock_pill"):
    rect(slide, l, t, w, h, fill=fill, line=line, rounded=True, adj=0.35, name=name)
    text(slide, l, t + 0.03, w, h, label, size=size, color=fg, bold=bold,
         align=PP_ALIGN.CENTER, name=name + "_t")


CARDS = [
    {
        "img": "mvp-anarkali.png",
        "brand": "Libas",
        "title": "Embroidered Anarkali Suit",
        "price": "Rs. 2499",
        "chip": "Fit confidence 67",
        "thin": False,
    },
    {
        "img": "mvp-lehenga.png",
        "brand": "Kalini",
        "title": "Embroidered Lehenga Set",
        "price": "Rs. 4499",
        "chip": "Need more reviews",
        "thin": True,
    },
]


def build_slide_7(slide):
    clear(slide)
    rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)

    text(slide, 0.5, 0.35, 12.33, 0.3, "MVP — FUNCTIONAL PROTOTYPE",
         size=11, bold=True, color=AMBER_DK)
    text(slide, 0.5, 0.62, 12.33, 0.9,
         "The MVP lives on the wishlist: check fit before Move to Bag",
         size=22, bold=True, color=NAVY)
    text(slide, 0.5, 1.48, 12.33, 0.35,
         "Same surface users already open — one extra action at the moment of hesitation, no new destination and no discount.",
         size=13, color=MUTED)

    # ---------------- phone: Myntra-style wishlist ----------------
    rect(slide, 0.5, 1.9, 3.75, 4.25, fill=BEZEL, rounded=True, adj=0.09, name="mock_bezel")
    rect(slide, 0.62, 2.02, 3.51, 4.01, fill=WHITE, rounded=True, adj=0.05, name="mock_screen")

    text(slide, 0.74, 2.08, 1.8, 0.28, "Wishlist", size=13, bold=True, color=NAVY, name="mock_h1")
    text(slide, 2.4, 2.12, 1.6, 0.24, "7 items", size=9.5, color=MUTED,
         align=PP_ALIGN.RIGHT, name="mock_count")
    rect(slide, 0.74, 2.42, 3.27, 0.012, fill=LINE, name="mock_rule")

    for idx, card in enumerate(CARDS):
        x = 0.74 + idx * 1.72
        img = IMG_DIR / card["img"]
        if img.exists():
            slide.shapes.add_picture(str(img), Inches(x + 0.05), Inches(2.52),
                                     Inches(1.45), Inches(1.45)).name = f"mock_img{idx}"
        else:
            rect(slide, x + 0.05, 2.52, 1.45, 1.45, fill=LIGHT, name=f"mock_img{idx}")
        text(slide, x + 1.22, 2.56, 0.26, 0.24, "♥", size=11, bold=True,
             color=PINK, align=PP_ALIGN.CENTER, name=f"mock_heart{idx}")
        pill(slide, x + 0.05, 3.66, 1.38, 0.26, card["chip"],
             fill=WHITE, fg=PINK if card["thin"] else NAVY, size=8,
             name=f"mock_chip{idx}")

        text(slide, x, 4.02, 1.6, 0.22, card["brand"], size=10, bold=True,
             color=NAVY, name=f"mock_brand{idx}")
        text(slide, x, 4.22, 1.6, 0.32, card["title"], size=8.5, color=MUTED,
             name=f"mock_title{idx}")
        text(slide, x, 4.54, 1.6, 0.22, card["price"], size=10, bold=True,
             color=INK, name=f"mock_price{idx}")

        pill(slide, x, 4.8, 1.55, 0.29, "MOVE TO BAG", fill=WHITE, fg=PINK,
             size=8, line=PINK, name=f"mock_bag{idx}")
        pill(slide, x, 5.13, 1.55, 0.29, "CHECK FIT", fill=PINK, fg=WHITE,
             size=8, name=f"mock_fit{idx}")

    text(slide, 0.74, 5.52, 3.27, 0.24, "Scroll for 5 more saved items",
         size=8.5, color=MUTED, align=PP_ALIGN.CENTER, name="mock_scroll")
    text(slide, 0.74, 5.75, 3.27, 0.22,
         "Prototype, not the live Myntra app", size=8.5, italic=True,
         color=MUTED, align=PP_ALIGN.CENTER, name="mock_note")

    # ---------------- assistant sheet ----------------
    rect(slide, 4.45, 1.9, 4.15, 4.25, fill=LIGHT, rounded=True, adj=0.04)
    text(slide, 4.65, 2.02, 3.75, 0.28, "CHECK FIT — OPENS ON THAT SAVED ITEM",
         size=10, bold=True, color=AMBER_DK)
    text(slide, 4.65, 2.3, 3.75, 0.26, "Libas · Embroidered Anarkali Suit",
         size=12.5, bold=True, color=NAVY)
    text(slide, 4.65, 2.54, 3.75, 0.24, "Saved 3 weeks · intent: Probably",
         size=10.5, color=MUTED)

    rect(slide, 4.65, 2.82, 3.75, 0.4, fill=AMBER, rounded=True, adj=0.12)
    text(slide, 4.65, 2.9, 3.75, 0.3, "Fit confidence 67 — labelled, not a black box",
         size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    text(slide, 4.65, 3.3, 3.75, 0.6,
         "From this item's 15 reviews: 3 read true-to-size, 3 say size up, 3 praise the threadwork, 2 flag average lining.",
         size=11, color=INK)

    rect(slide, 4.65, 3.94, 3.75, 0.52, fill=WHITE, rounded=True, adj=0.1, line=LINE)
    text(slide, 4.78, 4.0, 3.5, 0.44,
         "\u201cWill the work look cheap in person? Should I size up for daytime Rakhi?\u201d",
         size=10.5, italic=True, color=INK)

    text(slide, 4.65, 4.54, 3.75, 0.78,
         "The answer quotes those reviews back with counts, and says what the reviews do not cover — instead of inventing a percentage.",
         size=11, color=INK)

    pill(slide, 4.65, 5.38, 1.8, 0.32, "All 15 reviews shown", fill=WHITE,
         fg=NAVY, size=9.5, line=LINE, name="mock_prov")
    pill(slide, 6.6, 5.38, 1.8, 0.32, "Compare max 3 saved", fill=WHITE,
         fg=NAVY, size=9.5, line=LINE, name="mock_cmp")

    # ---------------- capabilities ----------------
    text(slide, 8.78, 2.02, 4.05, 0.28, "WHAT IT DOES", size=10, bold=True, color=AMBER_DK)
    caps = [
        ("Confidence score with provenance", "Counted from that item's reviews and return reasons"),
        ("Ask on the saved item", "Fit, fabric or threadwork versus the photo, occasion"),
        ("Personalised fit", "Cross-references the shopper's past orders where available"),
        ("Compare, capped at 3", "Only similar items already on the wishlist"),
    ]
    y = 2.32
    for i, (head, sub) in enumerate(caps, 1):
        rect(slide, 8.78, y, 0.28, 0.28, fill=AMBER, rounded=True, adj=0.2)
        text(slide, 8.78, y + 0.02, 0.28, 0.26, str(i), size=11, bold=True,
             color=NAVY, align=PP_ALIGN.CENTER)
        text(slide, 9.16, y - 0.01, 3.67, 0.26, head, size=11.5, bold=True, color=NAVY)
        text(slide, 9.16, y + 0.24, 3.67, 0.3, sub, size=10, color=MUTED)
        y += 0.66

    rect(slide, 8.78, 5.02, 4.05, 1.13, fill=LIGHT, rounded=True, adj=0.06)
    text(slide, 8.96, 5.14, 3.7, 0.26, "COLD START, HANDLED HONESTLY",
         size=10, bold=True, color=NAVY)
    text(slide, 8.96, 5.4, 3.7, 0.7,
         "The Kalini lehenga has 3 reviews, so its card reads \u201cNeed more reviews\u201d and falls back to brand size chart and Q&A.",
         size=10.5, color=MUTED)

    # ---------------- closing strip + link ----------------
    rect(slide, 0.5, 6.28, 12.33, 0.44, fill=NAVY)
    text(slide, 0.7, 6.35, 11.9, 0.32,
         "Moves stage 2 (decision confidence): Move to Bag stays the shopper's decision, and nothing in the flow offers a discount.",
         size=11.5, bold=True, color=AMBER)
    text_mixed(slide, 0.5, 6.78, 9.5, 0.26,
               [("Deployed MVP: ", {"bold": True, "color": NAVY}),
                ("[LINK: Deployed MVP — Streamlit Tab 2]", {"color": AMBER_DK})],
               size=11)

    text(slide, 0.5, 7.1, 6.0, 0.3, "Product Manager, Growth Team", size=10, color=MUTED)
    text(slide, 12.33, 7.1, 0.5, 0.3, "7", size=10, color=MUTED)


def normalize_font14(prs, minimum=14.0):
    changed = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name.startswith("mock_"):
                continue
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size and run.font.size.pt < minimum:
                        run.font.size = Pt(minimum)
                        changed += 1
    return changed


def main():
    prs = Presentation(str(SRC))
    build_slide_7(prs.slides[6])
    if "--font14" in sys.argv:
        n = normalize_font14(prs)
        print(f"bumped {n} runs to 14pt (phone mock left as UI chrome)")
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT}")


if __name__ == "__main__":
    main()
